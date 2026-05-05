
import re
import threading
import traceback
from datetime import datetime
from pathlib import Path
from tkinter import Tk, Label, Entry, Button, filedialog, Text, END, StringVar, messagebox

import pythoncom
import win32com.client
from pptx import Presentation


POWERPOINT_EXTENSIONS = {".pptx", ".pptm", ".ppt"}


# ---------- Dateiname bereinigen ----------
def safe_filename(name: str, fallback: str, max_length: int = 80) -> str:
    if not name or not str(name).strip():
        return fallback

    name = str(name)

    name = re.sub(r"[\x00-\x1f\x7f]", " ", name)
    name = re.sub(r'[<>:"/\\|?*]', "_", name)
    name = re.sub(r"\s+", " ", name)
    name = name.strip(" ._")

    name = name[:max_length]

    return name if name else fallback


# ---------- Dateien finden ----------
def find_powerpoints(source_dir: Path):
    for file in source_dir.rglob("*"):
        if file.suffix.lower() in POWERPOINT_EXTENSIONS:
            yield file


# ---------- Titel ----------
def get_slide_title(slide, index):
    if slide.shapes.title and slide.shapes.title.text.strip():
        return slide.shapes.title.text

    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            return shape.text.splitlines()[0]

    return f"Folie_{index:03d}"


# ---------- Texte ----------
def extract_slide_text(slide):
    texts = []
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            texts.append(shape.text.strip())
    return "\n".join(texts)


def extract_slide_notes(slide):
    if not slide.has_notes_slide:
        return ""

    texts = []
    for shape in slide.notes_slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            texts.append(shape.text.strip())
    return "\n".join(texts)


# ---------- Export ----------
def export_slides(powerpoint_file, temp_dir, log):
    pythoncom.CoInitialize()

    ppt = None
    pres = None

    try:
        ppt = win32com.client.DispatchEx("PowerPoint.Application")
        ppt.Visible = False

        pres = ppt.Presentations.Open(str(powerpoint_file), WithWindow=False)

        temp_dir.mkdir(parents=True, exist_ok=True)

        pres.Export(str(temp_dir), "PNG", 1920, 1080)

        # rekursiv + eindeutig
        pngs = list({
            p.resolve(): p
            for p in temp_dir.rglob("*")
            if p.suffix.lower() == ".png"
        }.values())

        pngs.sort(key=lambda p: int(re.search(r"\d+", p.name).group()))

        return pngs

    finally:
        if pres:
            pres.Close()
        if ppt:
            ppt.Quit()
        pythoncom.CoUninitialize()


# ---------- Hauptlogik ----------
def process_powerpoint(file, target_root, log):
    timestamp = datetime.now().strftime("%Y%m%d%H%M%S")
    result_dir = target_root / f"{timestamp}_{safe_filename(file.stem, 'ppt')}"
    result_dir.mkdir(parents=True, exist_ok=True)

    log(f"\nVerarbeite: {file}")

    prs = Presentation(file)

    titles = []
    texts_all = []
    notes_all = []

    for i, slide in enumerate(prs.slides, start=1):
        raw = get_slide_title(slide, i)
        clean = safe_filename(raw, f"Folie_{i:03d}")

        titles.append(clean)

        texts_all.append(f"{i:03d} {clean}\n{extract_slide_text(slide)}")
        notes = extract_slide_notes(slide)

        if notes:
            notes_all.append(f"{i:03d} {clean}\n{notes}")

    (result_dir / "Alle_Folientexte.txt").write_text("\n\n".join(texts_all), encoding="utf-8")
    (result_dir / "Alle_Notizen.txt").write_text("\n\n".join(notes_all), encoding="utf-8")

    temp = result_dir / "_tmp"
    pngs = export_slides(file, temp, log)

    for i, png in enumerate(pngs, start=1):
        title = titles[i - 1] if i <= len(titles) else f"Folie_{i:03d}"
        name = f"{i:03d}_{title}.png"
        png.replace(result_dir / name)

    # temp löschen
    for f in temp.rglob("*"):
        if f.is_file():
            f.unlink()
    for d in sorted(temp.rglob("*"), reverse=True):
        if d.is_dir():
            d.rmdir()
    temp.rmdir()

    log(f"Fertig: {file.name}")


# ---------- GUI ----------
class GUI:
    def __init__(self, root):
        self.root = root
        self.root.title("PowerPoint Extractor")

        self.src = StringVar()
        self.dst = StringVar()

        Label(root, text="Source").grid(row=0, column=0)
        Entry(root, textvariable=self.src, width=80).grid(row=0, column=1)
        Button(root, text="...", command=self.pick_src).grid(row=0, column=2)

        Label(root, text="Target").grid(row=1, column=0)
        Entry(root, textvariable=self.dst, width=80).grid(row=1, column=1)
        Button(root, text="...", command=self.pick_dst).grid(row=1, column=2)

        Button(root, text="Start", command=self.start).grid(row=2, column=1)

        self.logbox = Text(root, width=100, height=25)
        self.logbox.grid(row=3, column=0, columnspan=3)

    def pick_src(self):
        self.src.set(filedialog.askdirectory())

    def pick_dst(self):
        self.dst.set(filedialog.askdirectory())

    def log(self, msg):
        self.logbox.insert(END, msg + "\n")
        self.logbox.see(END)

    def start(self):
        threading.Thread(target=self.run).start()

    def run(self):
        src = Path(self.src.get())
        dst = Path(self.dst.get())

        for f in find_powerpoints(src):
            try:
                process_powerpoint(f, dst, self.log)
            except Exception as e:
                self.log(f"Fehler: {e}")
                self.log(traceback.format_exc())


if __name__ == "__main__":
    root = Tk()
    GUI(root)
    root.mainloop()
